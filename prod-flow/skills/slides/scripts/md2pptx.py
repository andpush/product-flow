#!/usr/bin/env python3
"""
Universal Markdown to PowerPoint Converter

Converts markdown files with --- slide separators into PPTX presentations.
Supports speaker notes, tables, bullets, links, inline formatting, and
multi-column layouts (2-3 columns using || delimiter).
"""

import re
import argparse
import yaml
from markdown_it import MarkdownIt
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor

# ============================================================================
# CONFIGURATION
# ============================================================================

DEFAULT_FONT = "Heebo"
ACCENT_COLOR = RGBColor(0x36, 0x4F, 0x6B)  # Smoky blue for titles, bold, italic, links
TEXT_COLOR = RGBColor(0, 0, 0)              # Black for regular text
MD_PARSER = MarkdownIt('commonmark', {'breaks': True, 'html': True}).enable('table')

# ============================================================================
# PARSING FUNCTIONS
# ============================================================================

def extract_frontmatter(markdown_content):
    """Extract YAML frontmatter. Returns: (frontmatter_dict, content_without_frontmatter)"""
    match = re.match(r'^---\s*\n(.*?)\n---\s*\n', markdown_content, re.DOTALL)
    if match:
        try:
            frontmatter = yaml.safe_load(match.group(1))
            return frontmatter, markdown_content[match.end():]
        except yaml.YAMLError:
            pass
    return None, markdown_content


def split_slides(markdown_content):
    """Split markdown content on --- separators"""
    return re.split(r'\n---\n', markdown_content)


def extract_notes(slide_content):
    """Extract speaker notes from HTML comments: <!-- Notes: text -->"""
    pattern = r'<!--\s*Notes:\s*(.*?)\s*-->'
    notes_match = re.search(pattern, slide_content, re.DOTALL)
    content = re.sub(pattern, '', slide_content, flags=re.DOTALL)
    notes = notes_match.group(1).strip() if notes_match else ""
    return content.strip(), notes


def parse_slide_content(content):
    """Parse markdown into title and body tokens."""
    tokens = MD_PARSER.parse(content)
    title = ""
    body_tokens = tokens

    # Find first h1 or h2 as title
    for i, token in enumerate(tokens):
        if token.type == 'heading_open' and token.tag in ['h1', 'h2']:
            if tokens[i + 1].type == 'inline':
                title = tokens[i + 1].content.strip()
                # Find heading_close and slice
                end_idx = i
                while tokens[end_idx].type != 'heading_close':
                    end_idx += 1
                body_tokens = tokens[end_idx + 1:]
                break

    return title, body_tokens


def parse_columns(raw_content):
    """Detect and parse column layout using || delimiter."""
    if '||' not in raw_content:
        return None

    # Remove title line before splitting
    lines = raw_content.split('\n')
    body_lines = []
    found_title = False
    for line in lines:
        if not found_title and line.strip().startswith('#'):
            found_title = True
            continue
        body_lines.append(line)

    body_content = '\n'.join(body_lines)
    if '||' not in body_content:
        return None

    columns = [col.strip() for col in body_content.split('||') if col.strip()]
    return columns if 2 <= len(columns) <= 3 else None


def determine_layout(title, body_tokens):
    """Determine slide layout: 1 (content) or 2 (section header)"""
    # Section header: has title but no body content
    if not body_tokens or all(t.type in ['paragraph_open', 'paragraph_close'] and
                               (t.type != 'inline' or not t.content.strip())
                               for t in body_tokens):
        has_content = any(t.type == 'inline' and t.content.strip() for t in body_tokens)
        if not has_content:
            return 2  # Section header
    return 1  # Content slide


# ============================================================================
# CONTENT RENDERING
# ============================================================================

def set_margins(text_frame):
    """Set reduced margins for text frames."""
    text_frame.margin_left = Inches(0.05)
    text_frame.margin_right = Inches(0.05)
    text_frame.margin_top = Inches(0.025)
    text_frame.margin_bottom = Inches(0.025)


def apply_inline_formatting(paragraph, inline_token):
    """Apply markdown inline formatting from markdown-it tokens."""
    if not inline_token.children:
        paragraph.text = inline_token.content
        return

    style_stack = []
    link_url = None

    for child in inline_token.children:
        if child.type in ['softbreak', 'hardbreak']:
            run = paragraph.add_run()
            run.text = '\n'
            continue

        # Handle style markers
        if child.type == 'strong_open':
            style_stack.append('bold')
            continue
        elif child.type == 'em_open':
            style_stack.append('italic')
            continue
        elif child.type == 'strong_close' and 'bold' in style_stack:
            style_stack.remove('bold')
            continue
        elif child.type == 'em_close' and 'italic' in style_stack:
            style_stack.remove('italic')
            continue
        elif child.type == 'link_open':
            if child.attrs and 'href' in child.attrs:
                link_url = child.attrs['href']
            style_stack.append('link')
            continue
        elif child.type == 'link_close' and 'link' in style_stack:
            style_stack.remove('link')
            link_url = None
            continue

        # Create run for content
        run = paragraph.add_run()
        run.text = child.content
        run.font.name = DEFAULT_FONT
        run.font.size = Pt(18)
        run.font.color.rgb = TEXT_COLOR

        if child.type == 'code_inline':
            run.font.name = 'Courier New'
            run.font.size = Pt(14)

        if 'bold' in style_stack:
            run.font.bold = True
            run.font.color.rgb = ACCENT_COLOR
        if 'italic' in style_stack:
            run.font.italic = True
            run.font.color.rgb = ACCENT_COLOR
        if 'link' in style_stack and link_url:
            run.font.color.rgb = ACCENT_COLOR
            run.font.underline = True
            try:
                run.hyperlink.address = link_url
            except:
                pass


def set_bullet(paragraph, level=0, enable=True):
    """Set or remove bullet formatting for a paragraph."""
    from pptx.oxml.ns import qn
    from lxml import etree

    pPr = paragraph._p.get_or_add_pPr()

    # Remove any existing bullet elements first
    for tag in ['buNone', 'buChar', 'buAutoNum', 'buBlip']:
        for elem in pPr.findall(qn(f'a:{tag}')):
            pPr.remove(elem)

    if enable:
        # Add bullet character (•) - let PowerPoint/Google Slides handle spacing
        buChar = etree.SubElement(pPr, qn('a:buChar'))
        buChar.set('char', '•')
        # Set level only - spacing will use application defaults
        paragraph.level = level
    else:
        # Explicitly disable bullets
        etree.SubElement(pPr, qn('a:buNone'))
        paragraph.level = 0


def add_content_from_tokens(text_frame, tokens):
    """Add content to text frame from markdown-it tokens."""
    text_frame.word_wrap = True
    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    set_margins(text_frame)

    # Clear placeholder text
    if text_frame.paragraphs:
        text_frame.paragraphs[0].clear()

    first_para_used = False
    list_level = -1  # -1 means not in a list
    i = 0

    while i < len(tokens):
        token = tokens[i]

        if token.type == 'paragraph_open':
            # Get or create paragraph
            if not first_para_used and text_frame.paragraphs:
                p = text_frame.paragraphs[0]
                first_para_used = True
            else:
                p = text_frame.add_paragraph()

            # Set bullet formatting based on whether we're in a list
            if list_level >= 0:
                # In a list: enable bullet with proper level
                set_bullet(p, level=list_level, enable=True)
            else:
                # Not in a list: disable bullet (plain paragraph)
                set_bullet(p, enable=False)

            # Apply content
            if i + 1 < len(tokens) and tokens[i + 1].type == 'inline':
                apply_inline_formatting(p, tokens[i + 1])
                i += 1
            i += 1

        elif token.type in ['bullet_list_open', 'ordered_list_open']:
            list_level += 1
            i += 1
        elif token.type in ['bullet_list_close', 'ordered_list_close']:
            list_level -= 1
            i += 1
        else:
            i += 1


def parse_table_from_tokens(tokens):
    """Extract table data from tokens."""
    headers, rows = [], []
    in_header, in_body, current_row = False, False, []

    for token in tokens:
        if token.type == 'table_close':
            break
        elif token.type == 'thead_open':
            in_header = True
        elif token.type == 'thead_close':
            in_header = False
        elif token.type == 'tbody_open':
            in_body = True
        elif token.type == 'tbody_close':
            in_body = False
        elif token.type == 'tr_open':
            current_row = []
        elif token.type == 'tr_close':
            if in_header:
                headers = current_row
            elif in_body:
                rows.append(current_row)
        elif token.type == 'inline':
            current_row.append(token.content)

    return headers, rows


def add_table_shape(slide, headers, rows):
    """Add a table to the slide."""
    if not headers:
        return

    num_rows, num_cols = len(rows) + 1, len(headers)
    left, top = Inches(0.3), Inches(0.85)
    width, height = Inches(9.4), Inches(4.8)

    shape = slide.shapes.add_table(num_rows, num_cols, left, top, width, height)
    table = shape.table

    for c in range(num_cols):
        table.columns[c].width = int(width / num_cols)

    for r in range(num_rows):
        table.rows[r].height = Inches(0.4)

    # Header row
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            for run in p.runs:
                run.font.bold = True
                run.font.name = DEFAULT_FONT

    # Data rows
    for r, row_data in enumerate(rows):
        for c, cell_text in enumerate(row_data):
            if c < num_cols:
                cell = table.cell(r + 1, c)
                cell.text = cell_text
                for p in cell.text_frame.paragraphs:
                    p.alignment = MSO_ANCHOR.MIDDLE
                    for run in p.runs:
                        run.font.name = DEFAULT_FONT


def add_column_layout(slide, columns):
    """Add 2 or 3 column layout using textboxes."""
    num_cols = len(columns)
    content_left, content_top = Inches(0.3), Inches(0.85)
    content_width, content_height = Inches(9.4), Inches(4.8)
    gap = Inches(0.2)

    total_gap = gap * (num_cols - 1)
    col_width = (content_width - total_gap) / num_cols

    for i, col_content in enumerate(columns):
        col_left = content_left + i * (col_width + gap)
        textbox = slide.shapes.add_textbox(col_left, content_top, col_width, content_height)
        col_tokens = MD_PARSER.parse(col_content)
        add_content_from_tokens(textbox.text_frame, col_tokens)


# ============================================================================
# SLIDE CREATION
# ============================================================================

def remove_shape(slide, shape):
    """Remove a shape from slide."""
    shape.element.getparent().remove(shape.element)


def style_title(title_shape, title_text, is_section=False):
    """Apply consistent title styling."""
    title_shape.text = ""
    title_frame = title_shape.text_frame
    title_frame.clear()

    p = title_frame.paragraphs[0]
    p.text = title_text
    title_size = Pt(44) if is_section else Pt(32)
    p.font.size = title_size
    p.line_spacing = 1.0
    p.space_before = Pt(0)
    p.space_after = Pt(0)

    for run in p.runs:
        run.font.color.rgb = ACCENT_COLOR
        run.font.name = DEFAULT_FONT
        run.font.bold = True
        run.font.size = title_size

    title_frame.word_wrap = False
    set_margins(title_frame)

    return title_shape


def create_title_slide(prs, frontmatter):
    """Create title slide from YAML frontmatter."""
    slide = prs.slides.add_slide(prs.slide_layouts[0])

    # Title
    if frontmatter.get('title') and slide.shapes.title:
        title_shape = style_title(slide.shapes.title, frontmatter['title'], is_section=True)
        title_shape.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        title_shape.top = Inches(1.2)
        title_shape.left = Inches(0.5)
        title_shape.width = Inches(9.0)
        title_shape.height = Inches(1.0)

    # Subtitle area
    if len(slide.placeholders) > 1:
        subtitle_shape = slide.placeholders[1]
        subtitle_shape.text = ""
        subtitle_frame = subtitle_shape.text_frame
        subtitle_frame.clear()
        subtitle_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        set_margins(subtitle_frame)

        # Build subtitle: subtitle, author, date (no labels)
        parts = [frontmatter.get(k) for k in ['subtitle', 'author', 'date'] if frontmatter.get(k)]

        for i, part in enumerate(parts):
            p = subtitle_frame.paragraphs[0] if i == 0 else subtitle_frame.add_paragraph()
            p.text = str(part)
            for run in p.runs:
                run.font.name = DEFAULT_FONT
                run.font.size = Pt(24)

        subtitle_shape.top = Inches(2.5)
        subtitle_shape.left = Inches(1.5)
        subtitle_shape.width = Inches(7.0)
        subtitle_shape.height = Inches(2.0)

    return slide


def create_section_slide(prs, title, notes=""):
    """Create a section header slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[2])

    if title and slide.shapes.title:
        title_shape = style_title(slide.shapes.title, title, is_section=True)
        title_shape.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        title_shape.top = Inches(1.2)
        title_shape.left = Inches(0.5)
        title_shape.width = Inches(9.0)
        title_shape.height = Inches(1.0)

    if notes:
        slide.notes_slide.notes_text_frame.text = notes

    return slide


def create_content_slide(prs, title, body_tokens, notes="", raw_content=""):
    """Create a content slide with title and body."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])

    # Title
    if title and slide.shapes.title:
        title_shape = style_title(slide.shapes.title, title)
        title_shape.text_frame.vertical_anchor = MSO_ANCHOR.TOP
        title_shape.left = Inches(0.3)
        title_shape.width = Inches(9.4)
        title_shape.height = Inches(0.5)

    # Remove default content placeholder - we'll use textbox instead
    for shape in list(slide.placeholders):
        if shape.placeholder_format.idx in [1, 14]:
            remove_shape(slide, shape)

    # Check for special content types
    columns = parse_columns(raw_content) if raw_content else None
    is_table = any(t.type == 'table_open' for t in body_tokens)

    if columns:
        add_column_layout(slide, columns)
    elif is_table:
        headers, rows = parse_table_from_tokens(body_tokens)
        add_table_shape(slide, headers, rows)
    elif body_tokens:
        # Use textbox for content (no default bullets)
        textbox = slide.shapes.add_textbox(Inches(0.3), Inches(0.85), Inches(9.4), Inches(4.8))
        add_content_from_tokens(textbox.text_frame, body_tokens)

    if notes:
        slide.notes_slide.notes_text_frame.text = notes

    return slide


# ============================================================================
# MAIN CONVERTER
# ============================================================================

def convert_markdown_to_pptx(markdown_file, output_file):
    """Main conversion function"""
    print(f"Reading markdown file: {markdown_file}")
    try:
        with open(markdown_file, 'r', encoding='utf-8') as f:
            markdown_content = f.read()
    except Exception as e:
        print(f"✗ Error reading file: {e}")
        return False

    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(10), Inches(5.625)

    # Extract frontmatter
    frontmatter, markdown_content = extract_frontmatter(markdown_content)
    slides_raw = split_slides(markdown_content)
    print(f"Found {len(slides_raw)} slides")

    # Title slide from frontmatter
    slide_num = 1
    if frontmatter:
        print(f"  Slide {slide_num}: Title - (from frontmatter)")
        create_title_slide(prs, frontmatter)
        slide_num += 1

    # Process remaining slides
    for slide_content in slides_raw:
        if not slide_content.strip():
            continue
        try:
            content, notes = extract_notes(slide_content)
            title, body_tokens = parse_slide_content(content)
            layout = determine_layout(title, body_tokens)

            layout_name = "Section" if layout == 2 else "Content"
            print(f"  Slide {slide_num}: {layout_name} - '{title[:50]}...'")

            if layout == 2:
                create_section_slide(prs, title, notes)
            else:
                create_content_slide(prs, title, body_tokens, notes, raw_content=content)

            slide_num += 1
        except Exception as e:
            print(f"⚠ Warning: Error processing slide {slide_num}: {e}")
            import traceback
            traceback.print_exc()
            continue

    try:
        prs.save(output_file)
        print(f"✓ Presentation saved: {output_file}")
        print(f"✓ Total slides: {len(prs.slides)}")
        return True
    except Exception as e:
        print(f"✗ Error saving file: {e}")
        return False


# ============================================================================
# CLI INTERFACE
# ============================================================================

def main():
    """CLI entry point"""
    parser = argparse.ArgumentParser(
        description='Convert markdown to PowerPoint presentation',
        formatter_class=argparse.RawDescriptionHelpFormatter)
    parser.add_argument('input', help='Input markdown file')
    parser.add_argument('-o', '--output', required=True, help='Output PPTX file')
    args = parser.parse_args()
    success = convert_markdown_to_pptx(args.input, args.output)
    exit(0 if success else 1)


if __name__ == "__main__":
    main()
